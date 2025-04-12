import streamlit as st
import os
import tempfile
from datetime import datetime
from dotenv import load_dotenv
from langchain.schema import HumanMessage, AIMessage
from googletrans import Translator

load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_API_KEY2 = os.getenv("GOOGLE_API_KEY2")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID")

# NLP + ML Libraries
import nltk
import numpy as np
from nltk.tokenize import sent_tokenize
from sklearn.cluster import KMeans
from sklearn.metrics.pairwise import cosine_similarity
from rag.langchain_tools import GoogleSearchTool

# File Parsing Libraries
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

# Your Modules (custom)
from text_preprocessing.cleaner import clean_text
from embedding.embedder import get_embeddings
from search.vector_search import VectorSearchEngine
from summarization.summarizer import flan_summarize_with_query as summarize_text_with_query
from rag.rag_engine import generate_answer_with_rag

# Download tokenizer if missing
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")

# ---------------------------
# Setup
# ---------------------------
st.set_page_config(page_title="SUMMARY.AI", layout="wide")
pdf_folder_path = r"C:\Users\Asus\OneDrive\Desktop\AI Model\pdf_files"
os.makedirs(pdf_folder_path, exist_ok=True)

def translate_text(text, dest_lang):
    # If text is None or empty, return an empty string
    if text is None or text == "":
        return ""
    try:
        translated = translator.translate(text, dest=dest_lang)
        return translated.text
    except Exception as e:
        return f"Translation Error: {e}"

translator = Translator()

# Language dropdown and mapping
lang_option = st.selectbox("Select Output Language", ["English", "Spanish", "French", "Hindi", "Chinese"])
output_lang_map = {
    "English": "en",
    "Spanish": "es",
    "French": "fr",
    "Hindi": "hi",
    "Chinese": "zh-cn"
}
selected_lang_code = output_lang_map.get(lang_option, "en")

# ---------------------------
# File Parsing Functions
# ---------------------------
def extract_text_from_file(file_path):
    ext = file_path.lower().split('.')[-1]
    if ext == 'pdf':
        reader = PdfReader(file_path)
        text = ""
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text += page_text + "\n"
            else:
                try:
                    with tempfile.TemporaryDirectory() as tmpdir:
                        images = convert_from_path(file_path, first_page=i+1, last_page=i+1, output_folder=tmpdir)
                        for img in images:
                            text += pytesseract.image_to_string(img) + "\n"
                except Exception as e:
                    print(f"OCR failed for page {i+1}: {e}")
        return text.strip()
    elif ext == 'docx':
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    elif ext == 'pptx':
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    elif ext == 'txt':
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        return ""

def get_documents_from_folder(folder_path):
    supported_exts = [".pdf", ".docx", ".pptx", ".txt"]
    documents = []
    for i, filename in enumerate(os.listdir(folder_path)):
        if not any(filename.lower().endswith(ext) for ext in supported_exts):
            continue
        path = os.path.join(folder_path, filename)
        text = extract_text_from_file(path)
        if len(text) < 100:
            continue
        documents.append({
            "doc_id": f"doc_{i:03}",
            "name": filename,
            "content": text,
            "metadata": {
                "author": "Unknown",
                "date": datetime.fromtimestamp(os.path.getmtime(path)).strftime('%Y-%m-%d')
            }
        })
    return documents

# ---------------------------
# Sidebar: File Upload
# ---------------------------
st.sidebar.header("📤 Upload Files")
uploaded_files = st.sidebar.file_uploader(
    "Drop PDF/DOCX/PPTX/TXT files here",
    accept_multiple_files=True,
    type=["pdf", "docx", "pptx", "txt"]
)
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_path = os.path.join(pdf_folder_path, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
    st.sidebar.success("✅ Files uploaded successfully! Please refresh the page to load them.")

# ---------------------------
# Load and Preprocess Documents
# ---------------------------
documents = get_documents_from_folder(pdf_folder_path)
if not documents:
    st.error("No valid documents found in the folder. Please upload or check content.")
    st.stop()

doc_texts = [clean_text(doc["content"]) for doc in documents]
doc_embeddings = get_embeddings(doc_texts)
search_engine = VectorSearchEngine(dimension=doc_embeddings.shape[1])
search_engine.build_index(doc_embeddings, doc_texts)
index_to_docinfo = {i: doc for i, doc in enumerate(documents)}

n_clusters = min(len(documents), 3)
topic_labels = KMeans(n_clusters=n_clusters, random_state=42).fit(doc_embeddings).labels_

# ---------------------------
# Fun & Vibrant UI Styling
# ---------------------------
st.markdown('<div class="big-title">📚 Smart AI-Powered Document Assistant</div>', unsafe_allow_html=True)
st.markdown("✨ **Explore, search, and chat with your files like never before!**")
with st.expander("💡 Tips for Fun Exploration"):
    st.write("""
        - Use the 🧠 Think mode to get deep insights from Google AI.
        - Search 🔍 mode gives you a precise summary from your documents.
        - Jump into **Conversation Mode** to ask follow-up questions; AI will let you know if it's off-topic.
        - Upload files on the sidebar and start your exploration!
    """)

lang_option = st.selectbox("🌐 Language (Coming Soon)", ["English", "Spanish", "French"])
voice_query = st.button("🎙 Try Voice Query (Coming Soon)")
if voice_query:
    st.info("🎤 Voice input is in the works — stay tuned!")

st.markdown("### 🧭 Choose Your Action")

# ---------------------------
# Main Query Input & Handling
# ---------------------------
query = st.text_input("💬 What do you want to know?")
col1, col2, col3 = st.columns(3)
with col1:
    search_button = st.button("🔎 Search")
with col2:
    think_button = st.button("🧠 Think")
with col3:
    explore_button = st.button("🌐 Explore More")

if "prev_queries" not in st.session_state:
    st.session_state["prev_queries"] = []
if "last_summary" not in st.session_state:
    st.session_state["last_summary"] = ""
if "last_docname" not in st.session_state:
    st.session_state["last_docname"] = ""

if query:
    st.session_state.prev_queries.append(query)
    clean_query = clean_text(query)
    query_embedding = get_embeddings([clean_query])

    # ---- SEARCH Flow: Query-Focused Summarization using Flan-T5 ----
    if search_button:
        st.subheader("🔍 Search Results (Query-Focused)")
        found_relevant = False
        relevant_results = []
        for doc in documents:
            sentences = sent_tokenize(doc["content"])
            if not sentences:
                continue
            sentence_embeddings = get_embeddings(sentences)
            sim_scores = np.dot(sentence_embeddings, query_embedding.T).flatten()
            best_score = np.max(sim_scores)
            best_idx = np.argmax(sim_scores)
            if best_score < 0.35:
                continue
            found_relevant = True
            start = max(0, best_idx - 1)
            end = min(len(sentences), best_idx + 2)
            relevant_chunk = " ".join(sentences[start:end])
            # Use Flan-T5 based summarization for a query-focused summary
            query_summary = generate_answer_with_rag(query, relevant_chunk, doc["name"])
            relevant_results.append({
                "doc": doc,
                "best_sentence": sentences[best_idx],
                "score": best_score,
                "summary": query_summary
            })
            # Save latest summary and doc name in session state for Think chain
            st.session_state["last_summary"] = query_summary
            st.session_state["last_docname"] = doc["name"]

        if not found_relevant:
            st.warning("🤷 No relevant documents found for your query.")
        else:
            for result in sorted(relevant_results, key=lambda x: x["score"], reverse=True):
                st.markdown(f"### 📄 {result['doc']['name']}")
                st.markdown(f"**💡 Best Matching Sentence:** {result['best_sentence']} (Score: {result['score']:.2f})")
                st.markdown(f"**📝 Summary:** {result['summary']}")

    # ---- THINK Flow: Use Google Generative API via LangChain for Deep Insight ----
    if think_button:
        st.subheader("🧠 Deep AI Insight")
        summary = st.session_state.get("last_summary", "")
        docname = st.session_state.get("last_docname", "")
        if not summary or not docname:
            st.warning("Please run a Search first to generate a summary.")
        else:
            try:
                from rag.Google_Gen_AI import generate_answer_with_google
                refined_answer = generate_answer_with_google(query, summary, docname)
                translated_answer = translate_text(refined_answer, selected_lang_code)
                st.markdown(translated_answer)
            except Exception as e:
                st.error(f"Error: {e}")

    # ---- Conversation Mode: Follow-Up using LangChain and Cosine Similarity ----
    if st.checkbox("Enter Conversation Mode"):
        st.subheader("Conversation Mode")
        memory_length = st.slider("Conversation Memory Length (Number of messages)", min_value=1, max_value=20, value=5)
        from langchain.memory import ConversationBufferMemory
        from langchain.chains import ConversationChain
        from langchain_google_genai.chat_models import ChatGoogleGenerativeAI

        memory = ConversationBufferMemory(return_messages=True)
        llm = ChatGoogleGenerativeAI(
            model="gemini-1.5-pro-001",
            temperature=0.7,
            google_api_key=GOOGLE_API_KEY
        )
        conversation_chain = ConversationChain(llm=llm, memory=memory)
        followup = st.text_input("Enter your follow-up question (or type 'done' to exit):")

        if st.button("Send Follow-Up"):
            if followup.lower() == "done":
                st.info("Ending conversation mode.")
                memory.clear()
            else:
                docname_context = st.session_state.get("last_docname", "")
                summary_context = st.session_state.get("last_summary", "")
                original_query = st.session_state.prev_queries[-1] if st.session_state.prev_queries else ""
                if not docname_context or not summary_context or not original_query:
                    st.warning("Please run Search or Think first to set the context.")
                else:
                    query_emb = get_embeddings([clean_text(original_query)])[0]
                    followup_emb = get_embeddings([clean_text(followup)])[0]
                    similarity = cosine_similarity([query_emb], [followup_emb])[0][0]
                    threshold = 0.5  # Tune this value as needed
                    if similarity >= threshold:
                        context_prompt = (
                            f"You are now engaged in a follow-up conversation based on the document titled '{docname_context}'.\n"
                            f"Below is a summary of that document:\n{summary_context}\n\n"
                            "When answering follow-up questions, provide detailed, accurate, and comprehensive responses. "
                            "If a particular factual detail (e.g., names, ages) is widely known but not mentioned in the summary, you may safely incorporate it from your verified general knowledge. "
                            "Ensure that you present the information clearly and accurately—avoid unnecessary repetition and be specific. "
                            "If you're confident in the answer (for example, if asked about Virat Kohli's wife and age), include those details explicitly. "
                            "However, if the context does not provide enough information and you are not certain, state that additional verified details are not available.\n\n"
                            f"Follow-up Question: {followup}"
                        )
                        convo_answer = conversation_chain.run(context_prompt)
                        translated_convo = translate_text(convo_answer, selected_lang_code)
                        st.markdown(f"🤖 AI:** {translated_convo}")
                    else:
                        st.warning("❗ Irrelevant question. Please ask something related to the document context.")

    # ---- EXPLORE Flow: Only Display Relevant Web Content (No Extra Insight) ----
    if explore_button:
        st.markdown("**🌍 Relevant Web Content Based on Document Insight:**")
        summary = st.session_state.get("last_summary", "")
        if summary:
            enhanced_query = f"{query} - Context: {summary[:250]}"
        else:
            enhanced_query = f"{query}"
        try:
            search_tool = GoogleSearchTool(api_key=GOOGLE_API_KEY2, cse_id=GOOGLE_CSE_ID)
            explore_results = search_tool.run(enhanced_query)
            st.markdown(explore_results)
        except Exception as e:
            st.error(f"Error fetching explore content: {e}")
    # ---- Previous Queries Display ----
    st.subheader("🗣 Previous Queries")
    for i, prev_q in enumerate(reversed(st.session_state.get("prev_queries", [])[:5]), 1):
        st.markdown(f"{i}. {prev_q}")
st.info("TXT support, drag & drop, and more intelligent clustering are live! 🎉 More coming soon!")
